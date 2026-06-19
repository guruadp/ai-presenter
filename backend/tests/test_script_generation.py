import io
import zipfile


def _sample_deck() -> bytes:
    try:
        from pptx import Presentation
    except ImportError:
        return _minimal_pptx()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Revenue"
    slide.placeholders[1].text = "ARR grew this quarter"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Close"
    slide.placeholders[1].text = "Summarize the launch story"

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()


def _minimal_pptx() -> bytes:
    stream = io.BytesIO()
    with zipfile.ZipFile(stream, "w") as zf:
        zf.writestr("[Content_Types].xml", "")
        zf.writestr(
            "ppt/slides/slide1.xml",
            """
            <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                   xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:cSld><p:spTree>
                <p:sp><p:txBody><a:p><a:r><a:t>Revenue</a:t></a:r></a:p></p:txBody></p:sp>
                <p:sp><p:txBody><a:p><a:r><a:t>ARR grew this quarter</a:t></a:r></a:p></p:txBody></p:sp>
              </p:spTree></p:cSld>
            </p:sld>
            """,
        )
        zf.writestr(
            "ppt/slides/slide2.xml",
            """
            <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                   xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:cSld><p:spTree>
                <p:sp><p:txBody><a:p><a:r><a:t>Close</a:t></a:r></a:p></p:txBody></p:sp>
                <p:sp><p:txBody><a:p><a:r><a:t>Summarize the launch story</a:t></a:r></a:p></p:txBody></p:sp>
              </p:spTree></p:cSld>
            </p:sld>
            """,
        )
    return stream.getvalue()


def test_generate_project_scripts_uses_slide_vision_tone_and_exact_kb_facts(client):
    kb = client.post("/kbs", json={"name": "Product KB", "owner": "alice"}).json()
    client.post(
        f"/kbs/{kb['id']}/facts",
        json={"key": "ARR", "value": "$4.2M", "source": "finance.md"},
    )
    project = client.post(
        "/projects",
        json={
            "name": "Launch deck",
            "owner": "alice",
            "kb_ids": [kb["id"]],
            "tone_profile": {
                "formality": "polished",
                "pace": "brisk",
                "persona": "product strategist",
                "dos": [],
                "donts": [],
                "language": "en",
                "voice_id": None,
            },
        },
    ).json()
    client.post(
        f"/projects/{project['id']}/deck",
        files={
            "file": (
                "launch.pptx",
                _sample_deck(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    resp = client.post(f"/projects/{project['id']}/scripts")

    assert resp.status_code == 200
    slides = resp.json()["slides"]
    assert len(slides) == 2
    first_script = slides[0]["script"]
    assert first_script["status"] == "draft"
    assert "$4.2M" in first_script["narration"]
    assert "Vision pass" not in first_script["narration"]
    assert "Rendered image" not in first_script["narration"]
    assert "the slide says" not in first_script["narration"].lower()
    assert first_script["duration_seconds"] > 0
    assert first_script["delivery_style"]["persona"] == "product strategist"
    assert first_script["delivery_style"]["pace"] == "brisk"
    assert first_script["segments"][0]["audio_tags"]
    assert first_script["citations"][0]["type"] == "structured_fact"
    assert first_script["citations"][0]["value"] == "$4.2M"
    assert first_script["citations"][0]["source"] == "finance.md"
    assert slides[1]["script"]["running_summary"]


def test_regenerate_single_slide_applies_feedback_without_touching_neighbor(client):
    project = client.post(
        "/projects",
        json={"name": "Launch deck", "owner": "alice"},
    ).json()
    project = client.post(
        f"/projects/{project['id']}/deck",
        files={
            "file": (
                "launch.pptx",
                _sample_deck(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    ).json()
    generated = client.post(f"/projects/{project['id']}/scripts").json()
    first_slide = generated["slides"][0]
    second_script = generated["slides"][1]["script"]

    resp = client.post(
        f"/projects/{project['id']}/slides/{first_slide['id']}/script/regenerate",
        json={
            "feedback": "Make this more executive and concise",
            "make_shorter": True,
            "more_energy": True,
        },
    )

    assert resp.status_code == 200
    regenerated = resp.json()
    assert regenerated["version"] == first_slide["script"]["version"] + 1
    assert regenerated["feedback"] == "Make this more executive and concise"
    assert regenerated["delivery_style"]["energy"] == "high"
    assert "Revision note applied" not in regenerated["narration"]
    assert "feedback" not in regenerated["narration"].lower()

    after = client.get(f"/projects/{project['id']}").json()
    assert after["slides"][1]["script"]["id"] == second_script["id"]
    assert after["slides"][1]["script"]["version"] == second_script["version"]


def test_generate_scripts_requires_uploaded_deck(client):
    project = client.post(
        "/projects",
        json={"name": "Empty", "owner": "alice"},
    ).json()

    resp = client.post(f"/projects/{project['id']}/scripts")

    assert resp.status_code == 400
    assert "Upload a deck" in resp.json()["detail"]
