import io
import os
import zipfile

from app.config import get_settings


def _sample_deck() -> bytes:
    try:
        from pptx import Presentation
    except ImportError:
        return _minimal_pptx()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Welcome"
    slide.placeholders[1].text = "Grounded presenter workflow"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Evidence"
    slide.placeholders[1].text = "Use pinned knowledge bases"

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()


def _minimal_pptx() -> bytes:
    stream = io.BytesIO()
    with zipfile.ZipFile(stream, "w") as zf:
        zf.writestr("[Content_Types].xml", "")
        zf.writestr(
            "ppt/slides/slide1.xml",
            """
            <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                   xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:cSld><p:spTree>
                <p:sp><p:txBody><a:p><a:r><a:t>Welcome</a:t></a:r></a:p></p:txBody></p:sp>
                <p:sp><p:txBody><a:p><a:r><a:t>Grounded presenter workflow</a:t></a:r></a:p></p:txBody></p:sp>
              </p:spTree></p:cSld>
            </p:sld>
            """,
        )
        zf.writestr(
            "ppt/slides/slide2.xml",
            """
            <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                   xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:cSld><p:spTree>
                <p:sp><p:txBody><a:p><a:r><a:t>Evidence</a:t></a:r></a:p></p:txBody></p:sp>
                <p:sp><p:txBody><a:p><a:r><a:t>Use pinned knowledge bases</a:t></a:r></a:p></p:txBody></p:sp>
              </p:spTree></p:cSld>
            </p:sld>
            """,
        )
    return stream.getvalue()


def test_create_project_pins_selected_kbs(client):
    kb = client.post("/kbs", json={"name": "Product KB", "owner": "alice"}).json()
    client.patch(f"/kbs/{kb['id']}", json={"name": "Product KB v2"})
    updated_kb = client.get(f"/kbs/{kb['id']}").json()

    resp = client.post(
        "/projects",
        json={
            "name": "Launch deck",
            "owner": "alice",
            "kb_ids": [kb["id"]],
            "tone_profile": {
                "formality": "polished",
                "pace": "brisk",
                "persona": "product strategist",
                "dos": ["cite concrete facts"],
                "donts": ["overpromise"],
                "language": "en",
                "voice_id": "voice_123",
            },
        },
    )

    assert resp.status_code == 201
    data = resp.json()
    assert data["name"] == "Launch deck"
    assert data["tone_profile"]["voice_id"] == "voice_123"
    assert data["knowledge_bases"][0]["kb_id"] == kb["id"]
    assert data["knowledge_bases"][0]["pinned_version"] == updated_kb["version"]
    assert data["knowledge_bases"][0]["pinned_content_hash"] == updated_kb["content_hash"]


def test_create_project_rejects_missing_kb(client):
    resp = client.post(
        "/projects",
        json={"name": "Bad project", "owner": "alice", "kb_ids": ["missing"]},
    )
    assert resp.status_code == 404


def test_upload_deck_parses_ordered_slides_and_images(client, tmp_path, monkeypatch):
    monkeypatch.setenv("STORAGE_DIR", str(tmp_path))
    get_settings.cache_clear()

    kb = client.post("/kbs", json={"name": "Product KB", "owner": "alice"}).json()
    project = client.post(
        "/projects",
        json={"name": "Launch deck", "owner": "alice", "kb_ids": [kb["id"]]},
    ).json()

    resp = client.post(
        f"/projects/{project['id']}/deck",
        files={
            "file": (
                "launch.pptx",
                _sample_deck(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert resp.status_code == 200
    slides = resp.json()["slides"]
    assert [slide["position"] for slide in slides] == [1, 2]
    assert slides[0]["title"] == "Welcome"
    assert "Grounded presenter workflow" in slides[0]["body"]
    assert slides[1]["title"] == "Evidence"
    assert all(slide["image_path"] for slide in slides)
    assert all(os.path.exists(slide["image_path"]) for slide in slides)

    get_settings.cache_clear()


def test_upload_deck_rejects_non_pptx(client):
    project = client.post("/projects", json={"name": "Launch deck", "owner": "alice"}).json()
    resp = client.post(
        f"/projects/{project['id']}/deck",
        files={"file": ("notes.txt", b"hello", "text/plain")},
    )
    assert resp.status_code == 400
    assert "Only .pptx" in resp.json()["detail"]
