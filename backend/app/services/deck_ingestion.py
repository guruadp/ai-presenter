import html
import io
import logging
import os
import re
import shutil
import struct
import subprocess
import tempfile
import zipfile
import zlib
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree

log = logging.getLogger(__name__)


@dataclass(frozen=True)
class ParsedSlide:
    position: int
    title: str | None
    body: str
    notes: str


def is_pptx(filename: str, content_type: str | None) -> bool:
    return filename.lower().endswith(".pptx") or content_type in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }


def parse_pptx(content: bytes) -> list[ParsedSlide]:
    try:
        return _parse_with_python_pptx(content)
    except ImportError:
        return _parse_with_xml(content)
    except Exception as e:
        raise ValueError("Invalid PPTX file") from e


def render_slide_images(
    slides: list[ParsedSlide],
    output_dir: str,
    pptx_path: str | None = None,
) -> list[str]:
    os.makedirs(output_dir, exist_ok=True)

    if pptx_path and os.path.exists(pptx_path):
        libreoffice_paths = _render_with_libreoffice(pptx_path, output_dir, slides)
        if libreoffice_paths:
            return libreoffice_paths

    image_paths: list[str] = []
    for slide in slides:
        path = Path(output_dir) / f"slide_{slide.position}.png"
        _render_preview_png(slide, path)
        image_paths.append(str(path))
    return image_paths


def _render_with_libreoffice(
    pptx_path: str,
    output_dir: str,
    slides: list[ParsedSlide],
) -> list[str]:
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
                capture_output=True,
                timeout=120,
            )
            if result.returncode != 0:
                log.warning("LibreOffice conversion failed: %s", result.stderr.decode())
                return []

            pdf_path = os.path.join(tmpdir, Path(pptx_path).stem + ".pdf")
            if not os.path.exists(pdf_path):
                log.warning("LibreOffice did not produce a PDF at %s", pdf_path)
                return []

            png_prefix = os.path.join(tmpdir, "slide")
            result = subprocess.run(
                ["pdftoppm", "-r", "150", "-png", pdf_path, png_prefix],
                capture_output=True,
                timeout=120,
            )
            if result.returncode != 0:
                log.warning("pdftoppm failed: %s", result.stderr.decode())
                return []

            image_paths: list[str] = []
            for slide in slides:
                src = _find_pdftoppm_output(tmpdir, slide.position)
                if src is None:
                    log.warning("No pdftoppm output for slide %d", slide.position)
                    return []
                dst = os.path.join(output_dir, f"slide_{slide.position}.png")
                shutil.copy2(src, dst)
                image_paths.append(dst)

            return image_paths
    except Exception:
        log.exception("LibreOffice slide rendering failed")
        return []


def _find_pdftoppm_output(tmpdir: str, position: int) -> str | None:
    for pad in range(1, 5):
        candidate = os.path.join(tmpdir, f"slide-{position:0{pad}d}.png")
        if os.path.exists(candidate):
            return candidate
    return None


def _parse_with_python_pptx(content: bytes) -> list[ParsedSlide]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise e

    prs = Presentation(io.BytesIO(content))
    slides: list[ParsedSlide] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        notes = ""
        if slide.has_notes_slide:
            notes = "\n".join(
                paragraph.text.strip()
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs
                if paragraph.text.strip()
            )
        title = texts[0] if texts else None
        body = "\n".join(texts[1:] if title else texts)
        slides.append(ParsedSlide(position=index, title=title, body=body, notes=notes))
    return slides


def _parse_with_xml(content: bytes) -> list[ParsedSlide]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(content))
    except zipfile.BadZipFile as e:
        raise ValueError("Invalid PPTX file") from e

    slide_names = sorted(
        (name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
        key=lambda name: int(re.search(r"slide(\d+)\.xml", name).group(1)),  # type: ignore[union-attr]
    )
    if not slide_names:
        raise ValueError("PPTX contains no slides")

    slides: list[ParsedSlide] = []
    for index, slide_name in enumerate(slide_names, start=1):
        texts = _extract_text_nodes(archive.read(slide_name))
        title = texts[0] if texts else None
        body = "\n".join(texts[1:] if title else texts)
        slides.append(ParsedSlide(position=index, title=title, body=body, notes=""))
    return slides


def _extract_text_nodes(xml_bytes: bytes) -> list[str]:
    root = ElementTree.fromstring(xml_bytes)
    texts = []
    for node in root.iter():
        if node.tag.endswith("}t") and node.text and node.text.strip():
            texts.append(html.unescape(node.text.strip()))
    return texts


def _render_preview_png(slide: ParsedSlide, path: Path) -> None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        _write_basic_png(path, slide)
        return

    width, height = 1280, 720
    image = Image.new("RGB", (width, height), color=(250, 252, 255))
    draw = ImageDraw.Draw(image)
    try:
        title_font = ImageFont.load_default(size=36)
        body_font = ImageFont.load_default(size=22)
        small_font = ImageFont.load_default(size=18)
    except TypeError:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()
        small_font = ImageFont.load_default()

    draw.rectangle((0, 0, width, 72), fill=(31, 41, 55))
    draw.text((36, 22), f"Slide {slide.position}", fill=(255, 255, 255), font=small_font)

    y = 115
    if slide.title:
        draw.multiline_text((64, y), _wrap(slide.title, 52), fill=(17, 24, 39), font=title_font, spacing=8)
        y += 100

    body = slide.body or "No body text extracted from this slide."
    draw.multiline_text((64, y), _wrap(body, 90), fill=(55, 65, 81), font=body_font, spacing=8)

    if slide.notes:
        draw.text((64, height - 64), "Speaker notes available", fill=(75, 85, 99), font=small_font)

    image.save(path, "PNG")


def _write_basic_png(path: Path, slide: ParsedSlide | None = None) -> None:
    width, height = 1280, 720
    pixels = bytearray((250, 252, 255) * width * height)

    _fill_rect(pixels, width, 0, 0, width, 72, (31, 41, 55))
    _draw_text(pixels, width, 36, 24, "SLIDE PREVIEW", (255, 255, 255), scale=3)

    title = slide.title if slide and slide.title else path.stem.replace("_", " ")
    body = slide.body if slide and slide.body else "No body text extracted from this slide."
    _draw_text(pixels, width, 64, 112, title, (17, 24, 39), scale=5)
    _draw_text(
        pixels,
        width,
        64,
        230,
        body[:500],
        (55, 65, 81),
        scale=3,
    )

    raw_rows = []
    for y in range(height):
        start = y * width * 3
        raw_rows.append(b"\x00" + bytes(pixels[start : start + width * 3]))
    raw = b"".join(raw_rows)

    def chunk(kind: bytes, data: bytes) -> bytes:
        checksum = zlib.crc32(kind + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + kind + data + struct.pack(">I", checksum)

    png = (
        b"\x89PNG\r\n\x1a\n"
        + chunk("IHDR".encode(), struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + chunk("IDAT".encode(), zlib.compress(raw))
        + chunk("IEND".encode(), b"")
    )
    path.write_bytes(png)


def _fill_rect(
    pixels: bytearray,
    width: int,
    x: int,
    y: int,
    rect_width: int,
    rect_height: int,
    color: tuple[int, int, int],
) -> None:
    height = len(pixels) // (width * 3)
    for py in range(max(0, y), min(height, y + rect_height)):
        for px in range(max(0, x), min(width, x + rect_width)):
            offset = (py * width + px) * 3
            pixels[offset : offset + 3] = bytes(color)


def _draw_text(
    pixels: bytearray,
    width: int,
    x: int,
    y: int,
    text: str,
    color: tuple[int, int, int],
    scale: int = 3,
) -> None:
    cursor_x = x
    cursor_y = y
    max_x = width - 64
    for char in text.upper():
        if char == "\n" or cursor_x + 6 * scale > max_x:
            cursor_x = x
            cursor_y += 9 * scale
            if char == "\n":
                continue
        if char == " ":
            cursor_x += 4 * scale
            continue
        glyph = _FONT.get(char, _FONT["?"])
        for row_index, row in enumerate(glyph):
            for col_index, bit in enumerate(row):
                if bit == "1":
                    _fill_rect(
                        pixels,
                        width,
                        cursor_x + col_index * scale,
                        cursor_y + row_index * scale,
                        scale,
                        scale,
                        color,
                    )
        cursor_x += 6 * scale


_FONT = {
    "A": ("01110", "10001", "10001", "11111", "10001", "10001", "10001"),
    "B": ("11110", "10001", "10001", "11110", "10001", "10001", "11110"),
    "C": ("01111", "10000", "10000", "10000", "10000", "10000", "01111"),
    "D": ("11110", "10001", "10001", "10001", "10001", "10001", "11110"),
    "E": ("11111", "10000", "10000", "11110", "10000", "10000", "11111"),
    "F": ("11111", "10000", "10000", "11110", "10000", "10000", "10000"),
    "G": ("01111", "10000", "10000", "10011", "10001", "10001", "01110"),
    "H": ("10001", "10001", "10001", "11111", "10001", "10001", "10001"),
    "I": ("11111", "00100", "00100", "00100", "00100", "00100", "11111"),
    "J": ("00111", "00010", "00010", "00010", "10010", "10010", "01100"),
    "K": ("10001", "10010", "10100", "11000", "10100", "10010", "10001"),
    "L": ("10000", "10000", "10000", "10000", "10000", "10000", "11111"),
    "M": ("10001", "11011", "10101", "10101", "10001", "10001", "10001"),
    "N": ("10001", "11001", "10101", "10011", "10001", "10001", "10001"),
    "O": ("01110", "10001", "10001", "10001", "10001", "10001", "01110"),
    "P": ("11110", "10001", "10001", "11110", "10000", "10000", "10000"),
    "Q": ("01110", "10001", "10001", "10001", "10101", "10010", "01101"),
    "R": ("11110", "10001", "10001", "11110", "10100", "10010", "10001"),
    "S": ("01111", "10000", "10000", "01110", "00001", "00001", "11110"),
    "T": ("11111", "00100", "00100", "00100", "00100", "00100", "00100"),
    "U": ("10001", "10001", "10001", "10001", "10001", "10001", "01110"),
    "V": ("10001", "10001", "10001", "10001", "10001", "01010", "00100"),
    "W": ("10001", "10001", "10001", "10101", "10101", "10101", "01010"),
    "X": ("10001", "10001", "01010", "00100", "01010", "10001", "10001"),
    "Y": ("10001", "10001", "01010", "00100", "00100", "00100", "00100"),
    "Z": ("11111", "00001", "00010", "00100", "01000", "10000", "11111"),
    "0": ("01110", "10001", "10011", "10101", "11001", "10001", "01110"),
    "1": ("00100", "01100", "00100", "00100", "00100", "00100", "01110"),
    "2": ("01110", "10001", "00001", "00010", "00100", "01000", "11111"),
    "3": ("11110", "00001", "00001", "01110", "00001", "00001", "11110"),
    "4": ("00010", "00110", "01010", "10010", "11111", "00010", "00010"),
    "5": ("11111", "10000", "10000", "11110", "00001", "00001", "11110"),
    "6": ("01110", "10000", "10000", "11110", "10001", "10001", "01110"),
    "7": ("11111", "00001", "00010", "00100", "01000", "01000", "01000"),
    "8": ("01110", "10001", "10001", "01110", "10001", "10001", "01110"),
    "9": ("01110", "10001", "10001", "01111", "00001", "00001", "01110"),
    ".": ("00000", "00000", "00000", "00000", "00000", "01100", "01100"),
    ",": ("00000", "00000", "00000", "00000", "01100", "00100", "01000"),
    ":": ("00000", "01100", "01100", "00000", "01100", "01100", "00000"),
    "-": ("00000", "00000", "00000", "11111", "00000", "00000", "00000"),
    "/": ("00001", "00010", "00010", "00100", "01000", "01000", "10000"),
    "$": ("00100", "01111", "10100", "01110", "00101", "11110", "00100"),
    "%": ("11001", "11010", "00010", "00100", "01000", "01011", "10011"),
    "?": ("01110", "10001", "00001", "00010", "00100", "00000", "00100"),
}


def _wrap(text: str, limit: int) -> str:
    lines: list[str] = []
    for raw_line in text.splitlines():
        words = raw_line.split()
        current: list[str] = []
        for word in words:
            if sum(len(part) for part in current) + len(current) + len(word) > limit:
                lines.append(" ".join(current))
                current = [word]
            else:
                current.append(word)
        if current:
            lines.append(" ".join(current))
    return "\n".join(lines)
